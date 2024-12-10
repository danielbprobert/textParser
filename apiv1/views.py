import time
from django.conf import settings
import uuid
import os
from rest_framework.views import APIView
from rest_framework.response import Response
from rest_framework import status
from PyPDF2 import PdfReader
from pdf2image import convert_from_path
from pytesseract import pytesseract, image_to_string
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from reportlab.pdfgen import canvas
from simple_salesforce import Salesforce
import requests
from transactions.models import Transaction, LogItem
from django.utils.timezone import now


# Set the path to the Tesseract executable (Windows specific)
pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe" 


class DocumentProcessingView(APIView):
    def post(self, request):
        transaction_id = str(uuid.uuid4())
        transaction = Transaction.objects.create(transaction_id=transaction_id, status="PENDING")
        transaction.start_time = now()
        transaction.save()

        start_time = time.time()

        fetch_parameters_log = LogItem.objects.create(
            transaction=transaction,
            step_name="Fetch Request Parameters",
            status="STARTED",
            start_time=now(),
        )
        session_id = request.data.get("sessionId")
        document_id = request.data.get("documentId")
        instance_url = request.data.get("instanceURL")

        if not session_id or not document_id or not instance_url:
            fetch_parameters_log.status = "FAILED"
            fetch_parameters_log.end_time = now()
            fetch_parameters_log.duration_ms = int((time.time() - start_time) * 1000)
            fetch_parameters_log.message = "Missing sessionId, documentId, or instanceURL"
            fetch_parameters_log.save()
            transaction.status = "FAILED"
            transaction.end_time = now()
            transaction.duration_ms = int((time.time() - start_time) * 1000)
            transaction.save()
            return Response(
                {"error": "Missing sessionId, documentId, or instanceURL", "transactionId": transaction_id},
                status=status.HTTP_400_BAD_REQUEST,
            )

        fetch_parameters_log.status = "COMPLETED"
        fetch_parameters_log.end_time = now()
        fetch_parameters_log.duration_ms = int((time.time() - start_time) * 1000)
        fetch_parameters_log.save()

        try:
            step_start_time = time.time()
            fetch_file_log = LogItem.objects.create(
                transaction=transaction,
                step_name="Fetch File from Salesforce",
                status="STARTED",
                start_time=now(),
            )
            file_path = self.fetch_file_from_salesforce(
                session_id, document_id, transaction_id, instance_url, fetch_file_log
            )
            fetch_file_log.duration_ms = int((time.time() - step_start_time) * 1000)
            fetch_file_log.save()

            step_start_time = time.time()
            convert_log = LogItem.objects.create(
                transaction=transaction,
                step_name="File Conversion",
                status="STARTED",
                start_time=now(),
            )

            file_extension = os.path.splitext(file_path)[1].lower()
            pdf_path = None

            if file_extension == ".pdf":
                pdf_path = file_path
            elif file_extension in [".jpg", ".jpeg", ".png"]:
                pdf_path = self.convert_image_to_pdf(file_path)
            elif file_extension == ".docx":
                pdf_path = self.convert_docx_to_pdf(file_path)
            elif file_extension in [".xls", ".xlsx"]:
                pdf_path = self.convert_excel_to_pdf(file_path)
            elif file_extension == ".pptx":
                pdf_path = self.convert_ppt_to_pdf(file_path)
            else:
                convert_log.status = "FAILED"
                convert_log.end_time = now()
                convert_log.duration_ms = int((time.time() - step_start_time) * 1000)
                convert_log.message = "Unsupported file type."
                convert_log.save()
                transaction.status = "FAILED"
                transaction.end_time = now()
                transaction.duration_ms = int((time.time() - start_time) * 1000)
                transaction.save()
                return Response(
                    {"error": "Unsupported file type", "transactionId": transaction_id},
                    status=status.HTTP_400_BAD_REQUEST,
                )

            convert_log.status = "COMPLETED"
            convert_log.end_time = now()
            convert_log.duration_ms = int((time.time() - step_start_time) * 1000)
            convert_log.save()

            step_start_time = time.time()
            extract_log = LogItem.objects.create(
                transaction=transaction,
                step_name="Extract Text from PDF",
                status="STARTED",
                start_time=now(),
            )

            parsed_text, num_pages, num_characters = self.extract_text_with_ocr(pdf_path)

            transaction.num_pages = num_pages
            transaction.num_characters = num_characters

            extract_log.status = "COMPLETED"
            extract_log.end_time = now()
            extract_log.duration_ms = int((time.time() - step_start_time) * 1000)
            extract_log.message = f"Extracted {num_characters} characters from {num_pages} pages."
            extract_log.save()

            transaction.status = "SUCCESS"
            transaction.end_time = now()
            transaction.duration_ms = int((time.time() - start_time) * 1000)
            transaction.save()

            return Response(
                {
                    "transactionId": transaction_id,
                    "numPages": num_pages,
                    "numCharacters": num_characters,
                    "durationMs": transaction.duration_ms,
                    "parsedText": parsed_text,
                },
                status=status.HTTP_200_OK,
            )

        except Exception as e:
            transaction.status = "FAILURE"
            transaction.end_time = now()
            transaction.duration_ms = int((time.time() - start_time) * 1000)
            transaction.save()

            error_log = LogItem.objects.create(
                transaction=transaction,
                step_name="Error Handling",
                status="FAILED",
                start_time=now(),
                end_time=now(),
                duration_ms=int((time.time() - start_time) * 1000),
                message=str(e),
            )
            return Response(
                {"error": str(e), "transactionId": transaction_id},
                status=status.HTTP_500_INTERNAL_SERVER_ERROR,
            )

    def extract_text_with_ocr(self, pdf_path):
        """
        Extract text using OCR for non-searchable PDFs.
        """
        reader = PdfReader(pdf_path)
        text = ""
        num_pages = len(reader.pages)

        for page in reader.pages:
            text += page.extract_text()

        # If text is empty, perform OCR
        if not text.strip():
            # Specify Poppler path for Windows
            poppler_path = r"C:\poppler\Library\bin"  # Update this path if necessary
            images = convert_from_path(pdf_path, dpi=300, poppler_path=poppler_path)
            ocr_text = [image_to_string(image, lang="eng") for image in images]
            text = " ".join(ocr_text)

        num_characters = len(text)
        return text, num_pages, num_characters

    def fetch_file_from_salesforce(self, session_id, document_id, transaction_id, instance_url, fetch_file_log):
        sf = Salesforce(instance_url=instance_url, session_id=session_id)

        content_version = sf.query(
            f"SELECT VersionData, Title, FileExtension FROM ContentVersion WHERE Id = '{document_id}'"
        )

        if not content_version["records"]:
            fetch_file_log.status = "FAILED"
            fetch_file_log.end_time = now()
            fetch_file_log.message = "No file found for provided DocumentId"
            fetch_file_log.save()
            raise ValueError(f"Transaction {transaction_id}: No file found for DocumentId {document_id}")

        version_data_relative_url = content_version["records"][0]["VersionData"]
        file_name = content_version["records"][0]["Title"]
        file_extension = content_version["records"][0]["FileExtension"]

        version_data_url = f"{instance_url}{version_data_relative_url}"

        headers = {"Authorization": f"Bearer {session_id}"}
        response = requests.get(version_data_url, headers=headers, stream=True)

        if response.status_code != 200:
            fetch_file_log.status = "FAILED"
            fetch_file_log.end_time = now()
            fetch_file_log.message = "Failed to connect to Salesforce"
            fetch_file_log.save()
            raise ValueError(f"Transaction {transaction_id}: Failed to fetch file content. HTTP Status {response.status_code}")

        file_path = os.path.join(settings.MEDIA_ROOT, f"{file_name}.{file_extension}")
        with open(file_path, "wb") as f:
            for chunk in response.iter_content(chunk_size=1024):
                f.write(chunk)

        fetch_file_log.status = "COMPLETED"
        fetch_file_log.end_time = now()
        fetch_file_log.save()
        return file_path

    def convert_image_to_pdf(self, image_path):
        pdf_path = f"{os.path.splitext(image_path)[0]}.pdf"
        images = convert_from_path(image_path)
        images[0].save(pdf_path, "PDF")
        return pdf_path

    def convert_docx_to_pdf(self, docx_path):
        pdf_path = f"{os.path.splitext(docx_path)[0]}.pdf"
        pdf = canvas.Canvas(pdf_path)
        document = Document(docx_path)
        text = "\n".join([p.text for p in document.paragraphs])
        pdf.drawString(100, 750, text)
        pdf.save()
        return pdf_path

    def convert_excel_to_pdf(self, excel_path):
        pdf_path = f"{os.path.splitext(excel_path)[0]}.pdf"
        workbook = load_workbook(excel_path)
        pdf = canvas.Canvas(pdf_path)

        y = 750
        margin = 50
        page_width = 595.27
        page_height = 841.89

        for sheet in workbook.sheetnames:
            worksheet = workbook[sheet]
            pdf.drawString(margin, y, f"Worksheet: {sheet}")
            y -= 20

            for row in worksheet.iter_rows(values_only=True):
                text = " | ".join([str(cell) if cell is not None else "" for cell in row])
                pdf.drawString(margin, y, text)
                y -= 20

                if y <= margin:
                    pdf.showPage()
                    y = page_height - margin

        pdf.save()
        return pdf_path

    def convert_ppt_to_pdf(self, ppt_path):
        pdf_path = f"{os.path.splitext(ppt_path)[0]}.pdf"
        presentation = Presentation(ppt_path)
        pdf = canvas.Canvas(pdf_path)
        y = 750
        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text
                    pdf.drawString(100, y, text)
                    y -= 20
        pdf.save()
        return pdf_path
